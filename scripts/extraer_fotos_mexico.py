"""
Extrae las fotos de los presidentes de colegio del PPT de investigación de México.

Cada diapositiva de colegio trae 4-5 imágenes: un fondo, una marca de agua y un
logo genérico que se repiten idénticos en todas (se detectan por hash), más el
logo del colegio (a la izquierda) y la foto de la persona (a la derecha). El
criterio posicional es estable en las 32: la foto de la persona siempre está en
la mitad derecha de la diapositiva.

Las que dicen "Sin Foto" en el PPT quedan sin archivo, a propósito.

Uso:  python scripts/extraer_fotos_mexico.py
Salida: scripts/fotos_mexico/<slug>.<ext>  +  scripts/fotos_mexico/_mapa.json
"""
import hashlib
import io
import json
import os
import re
import unicodedata
from collections import Counter

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

PPTX = "Presentación Investigacion UINL Mexico.pptx"
OUT = os.path.join("scripts", "fotos_mexico")
PICTURE = 13
ANCHO_DIAPO_IN = 13.33
# Se muestran como avatar de 40px (o 64px en la ficha). Guardarlas a tamaño
# original serían ~5 MB para la lista de México entera, y en el wifi de un
# congreso eso se paga caro. 640px en JPEG deja cada una en ~40 KB.
LADO_MAX = 640
CALIDAD = 85


def slug(s: str) -> str:
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode()
    s = re.sub(r"[^a-zA-Z0-9]+", "-", s).strip("-").lower()
    return s


def subdivision(titulo: str) -> str:
    """'COLEGIO DE NOTARIOS DEL ESTADO DE JALISCO' -> 'Jalisco'"""
    s = " ".join(titulo.split())
    if re.search(r"ciudad de m[eé]xico", s, re.I):
        return "Ciudad de México"
    if re.search(r"estado de m[eé]xico", s, re.I):
        return "Estado de México"
    if re.search(r"notariado mexicano", s, re.I):
        return "Colegio Nacional"
    s = re.sub(r"^colegio\s+de\s+notarios\s+p[uú]blicos?\s+", "", s, flags=re.I)
    s = re.sub(r"^colegio\s+de\s+notarios\s+", "", s, flags=re.I)
    s = re.sub(r"^consejo\s+de\s+notarios\s+", "", s, flags=re.I)
    s = re.sub(r"^(del|para el)\s+estado\s+", "", s, flags=re.I)
    s = re.sub(r"^estado\s+", "", s, flags=re.I)
    s = re.sub(r"^de\s+", "", s, flags=re.I)
    return s.strip().title()


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation(PPTX)
    slides = list(prs.slides)

    # 1 · Las imágenes decorativas son las que aparecen en muchas diapositivas.
    conteo = Counter()
    for slide in slides:
        vistos = set()
        for sh in slide.shapes:
            if sh.shape_type == PICTURE:
                h = hashlib.md5(sh.image.blob).hexdigest()
                if h not in vistos:
                    vistos.add(h)
                    conteo[h] += 1
    decorativas = {h for h, n in conteo.items() if n >= 5}

    mapa = []
    for i, slide in enumerate(slides, 1):
        if not (44 <= i <= 77):
            continue

        titulo = ""
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = " ".join(sh.text_frame.text.split())
                if re.search(r"colegio|consejo", t, re.I) and len(t) > len(titulo):
                    titulo = t
        if not titulo:
            continue

        sub = subdivision(titulo)
        sin_foto = any(
            sh.has_text_frame and re.search(r"sin\s+foto", sh.text_frame.text, re.I)
            for sh in slide.shapes
        )

        # 2 · Candidatas: no decorativas y en la mitad derecha.
        candidatas = []
        for sh in slide.shapes:
            if sh.shape_type != PICTURE:
                continue
            h = hashlib.md5(sh.image.blob).hexdigest()
            if h in decorativas:
                continue
            izq = Emu(sh.left).inches if sh.left else 0
            ancho = Emu(sh.width).inches if sh.width else 0
            if izq + ancho / 2 > ANCHO_DIAPO_IN / 2:
                candidatas.append((ancho * (Emu(sh.height).inches if sh.height else 0), sh))

        entrada = {"slide": i, "subdivision": sub, "titulo": titulo, "archivo": None}
        if candidatas and not sin_foto:
            _, sh = max(candidatas, key=lambda c: c[0])  # la más grande del lado derecho
            im = Image.open(io.BytesIO(sh.image.blob))
            if im.mode in ("RGBA", "P", "LA"):
                fondo = Image.new("RGB", im.size, (255, 255, 255))
                fondo.paste(im, mask=im.convert("RGBA").split()[-1])
                im = fondo
            im.thumbnail((LADO_MAX, LADO_MAX), Image.LANCZOS)
            nombre = f"{slug(sub)}.jpg"
            im.convert("RGB").save(os.path.join(OUT, nombre), "JPEG", quality=CALIDAD, optimize=True)
            entrada["archivo"] = nombre
        mapa.append(entrada)

    with open(os.path.join(OUT, "_mapa.json"), "w", encoding="utf-8") as f:
        json.dump(mapa, f, ensure_ascii=False, indent=2)

    con = sum(1 for m in mapa if m["archivo"])
    print(f"Diapositivas de colegio procesadas: {len(mapa)}")
    print(f"Con foto: {con}  ·  Sin foto: {len(mapa) - con}")
    for m in mapa:
        print(f"  slide {m['slide']:>3}  {m['subdivision']:<22} {m['archivo'] or '— sin foto —'}")


if __name__ == "__main__":
    main()
