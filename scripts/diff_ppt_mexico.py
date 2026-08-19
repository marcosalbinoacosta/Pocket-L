"""
Compara el PPT de investigación de México contra lo que ya está en la base.

Solo lee y reporta: no toca ni el PPT ni la base. Sirve para ver qué cambió
cuando llega una versión nueva de la presentación, antes de decidir importar.

Escanea todas las diapositivas buscando fichas de colegio en vez de asumir un
rango fijo, porque al actualizar el PPT se pueden agregar o mover slides.

Uso:  python scripts/diff_ppt_mexico.py
Salida: scripts/fotos_mexico/_ppt_actual.json
"""
import hashlib
import json
import os
import re
import unicodedata
from collections import Counter

from pptx import Presentation
from pptx.util import Emu

PPTX = "Presentación Investigacion UINL Mexico.pptx"
OUT = os.path.join("scripts", "fotos_mexico")
PICTURE = 13
ANCHO_DIAPO_IN = 13.33


def norm(s):
    s = unicodedata.normalize("NFKD", str(s or "")).encode("ascii", "ignore").decode()
    return re.sub(r"[^a-z0-9]", "", s.lower())


def subdivision(titulo: str) -> str:
    s = " ".join(titulo.split())
    if re.search(r"ciudad de m[eé]xico", s, re.I):
        return "Ciudad de México"
    if re.search(r"estado de m[eé]xico", s, re.I):
        return "Estado de México"
    if re.search(r"notariado mexicano", s, re.I):
        return "Colegio Nacional"
    s = re.sub(r"^colegio\s+de\s+notarios\s+p[uú]blicos?\s+", "", s, flags=re.I)
    s = re.sub(r"^colegio\s+de\s+notarios\s+", "", s, flags=re.I)
    s = re.sub(r"^consejo\s+de\s+notarios\s+", "", s, flags=re.I)
    s = re.sub(r"^(del|para el)\s+estado\s+", "", s, flags=re.I)
    s = re.sub(r"^estado\s+", "", s, flags=re.I)
    s = re.sub(r"^de\s+", "", s, flags=re.I)
    return s.strip()


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation(PPTX)
    slides = list(prs.slides)

    conteo = Counter()
    for slide in slides:
        for h in {hashlib.md5(sh.image.blob).hexdigest() for sh in slide.shapes if sh.shape_type == PICTURE}:
            conteo[h] += 1
    decorativas = {h for h, n in conteo.items() if n >= 5}

    fichas = []
    for i, slide in enumerate(slides, 1):
        textos = [" ".join(sh.text_frame.text.split()) for sh in slide.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()]
        crudos = [sh.text_frame.text.strip() for sh in slide.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()]

        # Una ficha de colegio es la que trae "Existen N Notarios"
        notarios = consumo = None
        for t in textos:
            m = re.search(r"Existen\s+([\d.,]+)\s+Notarios", t, re.I)
            if m:
                notarios = int(re.sub(r"[.,]", "", m.group(1)))
        for j, t in enumerate(crudos):
            if re.search(r"Consumo anual", t, re.I) and j + 1 < len(crudos):
                m2 = re.search(r"([\d.,]+)", crudos[j + 1])
                if m2:
                    consumo = int(re.sub(r"[.,]", "", m2.group(1)))
        if notarios is None and consumo is None:
            continue

        titulo = ""
        for t in textos:
            if re.search(r"colegio|consejo", t, re.I) and len(t) > len(titulo):
                titulo = t
        if not titulo:
            continue

        sin_foto = any(re.search(r"sin\s+foto", t, re.I) for t in textos)
        fotos = []
        for sh in slide.shapes:
            if sh.shape_type != PICTURE:
                continue
            h = hashlib.md5(sh.image.blob).hexdigest()
            if h in decorativas:
                continue
            izq = Emu(sh.left).inches if sh.left else 0
            anc = Emu(sh.width).inches if sh.width else 0
            alt = Emu(sh.height).inches if sh.height else 0
            if izq + anc / 2 > ANCHO_DIAPO_IN / 2:
                fotos.append((anc * alt, h))

        fichas.append({
            "slide": i,
            "titulo": titulo,
            "subdivision": subdivision(titulo),
            "notarios": notarios,
            "consumo": consumo,
            "hash_foto": (max(fotos)[1] if fotos and not sin_foto else None),
        })

    with open(os.path.join(OUT, "_ppt_actual.json"), "w", encoding="utf-8") as f:
        json.dump(fichas, f, ensure_ascii=False, indent=2)

    print(f"Fichas de colegio encontradas en el PPT: {len(fichas)}")
    for f_ in fichas:
        print(f"  slide {f_['slide']:>3}  {f_['subdivision']:<22} "
              f"notarios={str(f_['notarios']):<7} consumo={str(f_['consumo']):<10} "
              f"{'foto' if f_['hash_foto'] else 'SIN FOTO'}")


if __name__ == "__main__":
    main()
